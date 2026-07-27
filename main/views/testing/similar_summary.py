# Django에서 필요한 import
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from tempfile import NamedTemporaryFile

# 텍스트 추출 관련 라이브러리
import fitz  # PyMuPDF
from pptx import Presentation
from openpyxl import load_workbook

import os
import re
import json
from datetime import date
from main.request_logging import set_request_log_context
from main.utils.gemini_gemma import GemmaConfigError, GemmaGenerationError
from .similar_GPT import (
    generate_recommended_summaries,
    rerank_multiple_similar_candidates,
    run_gemini_gemma,
)
from .similar_compare import (
    SimilarSearchDependencyError,
    compare_multiple_from_index,
)


# PDF 파일에서 텍스트 추출
def parse_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text("text")
    return text


# DOCX 파일에서 텍스트 추출
def parse_docx(file_path):
    from zipfile import ZipFile
    from lxml import etree, objectify

    WORD_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    with ZipFile(file_path) as z:
        xml = z.read("word/document.xml")
    root = objectify.fromstring(xml)

    text_blocks = []

    for child in root.body.iterchildren():
        tag = child.tag.replace(WORD_NS, "")
        if tag == "p":  # 문단
            p_text = " ".join(t.text for t in child.iter(tag=WORD_NS+"t") if t.text)
            if p_text.strip():
                text_blocks.append(p_text.strip())
        elif tag == "tbl":  # 표
            for row in child.iter(tag=WORD_NS+"tr"):
                cells = []
                for tc in row.iter(tag=WORD_NS+"tc"):
                    tcPr = tc.tcPr if hasattr(tc, 'tcPr') else None
                    vmerge = None
                    if tcPr is not None and hasattr(tcPr, 'vMerge'):
                        vmerge = getattr(tcPr.vMerge, "val", None)
                        if vmerge is None or vmerge == "continue":
                            continue
                    cell_text = " ".join(t.text for t in tc.iter(tag=WORD_NS+"t") if t.text)
                    if cell_text.strip():
                        cells.append(cell_text.strip())
                if cells:
                    text_blocks.append(" | ".join(cells))

    txt = "\n".join(text_blocks)
    txt = re.sub(r'(\n\s*){2,}', '\n', txt)
    return txt.strip()


# PPTX 파일에서 텍스트 추출
def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text.extend([p.text for p in shape.text_frame.paragraphs])
    return "\n".join(text)


def parse_xlsx(file_path):
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    text_rows = []
    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows(values_only=True):
                values = [
                    str(value).strip()
                    for value in row
                    if value is not None and str(value).strip()
                ]
                if values:
                    text_rows.append(" | ".join(values))
    finally:
        workbook.close()
    return "\n".join(text_rows)


def parse_txt(file_path):
    for encoding in ("utf-8-sig", "cp949", "utf-8"):
        try:
            with open(file_path, "r", encoding=encoding) as text_file:
                return text_file.read()
        except UnicodeDecodeError:
            continue
    return None


# 파일 파싱 (Django UploadedFile 객체 활용)
def parse_file(uploaded_file):
    with NamedTemporaryFile(delete=False, suffix='.' + uploaded_file.name.split('.')[-1]) as tmp:
        for chunk in uploaded_file.chunks():
            tmp.write(chunk)
        tmp_path = tmp.name

    try:
        ext = uploaded_file.name.split('.')[-1].lower()
        if ext == 'pdf':
            return parse_pdf(tmp_path)
        elif ext == 'docx':
            return parse_docx(tmp_path)
        elif ext == 'pptx':
            return parse_pptx(tmp_path)
        elif ext == 'xlsx':
            return parse_xlsx(tmp_path)
        elif ext == 'txt':
            return parse_txt(tmp_path)
        else:
            return None
    finally:
        os.unlink(tmp_path)


# 텍스트 전처리 (공백 및 줄바꿈 제거)
def preprocess_text(text):
    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def _prepare_summary_options(request):
    file_type = request.POST.get("fileType", "")
    uploaded_file = request.FILES.get("file")
    manual_input = request.POST.get("manualInput", "").strip()

    if uploaded_file:
        set_request_log_context(
            request,
            feature="similar",
            input_mode="file",
            file_type=file_type,
            file_name=uploaded_file.name,
        )
        text = parse_file(uploaded_file)
        if text is None or len(text.strip()) < 10:
            return JsonResponse(
                {"response": "내용이 부족하거나 지원되지 않는 형식입니다."},
                status=400,
            )

        clean_text = preprocess_text(text)
        sentences = re.split(r"(?<=[.!?])\s+", clean_text)
        original_summary = preprocess_text(run_gemini_gemma(sentences))
        if original_summary.startswith("❌"):
            return JsonResponse({"response": original_summary}, status=503)

        try:
            recommendations = generate_recommended_summaries(
                clean_text,
                count=4,
                max_chars=60,
            )
        except (GemmaConfigError, GemmaGenerationError) as exc:
            return JsonResponse(
                {"response": f"추천 요약 문장을 생성하지 못했습니다: {exc}"},
                status=503,
            )

        options = [
            {
                "id": f"recommendation-{index + 1}",
                "text": summary,
                "is_original": False,
            }
            for index, summary in enumerate(recommendations)
        ]
        options.append(
            {
                "id": "original",
                "text": original_summary,
                "is_original": True,
            }
        )
        default_selected_ids = ["recommendation-1"]
        set_request_log_context(request, llm_summary=original_summary)
        input_mode = "file"
    elif manual_input:
        original_summary = preprocess_text(manual_input)
        set_request_log_context(
            request,
            feature="similar",
            input_mode="manual",
            manual_input=manual_input,
        )
        try:
            recommendations = generate_recommended_summaries(
                original_summary,
                count=5,
                max_chars=60,
            )
        except (GemmaConfigError, GemmaGenerationError) as exc:
            return JsonResponse(
                {"response": f"추천 제품 개요 문장을 생성하지 못했습니다: {exc}"},
                status=503,
            )

        options = [
            {
                "id": f"recommendation-{index + 1}",
                "text": summary,
                "is_original": False,
            }
            for index, summary in enumerate(recommendations)
        ]
        options.append(
            {
                "id": "original",
                "text": original_summary,
                "is_original": True,
            }
        )
        default_selected_ids = ["original"]
        input_mode = "manual"
    else:
        set_request_log_context(request, feature="similar", input_mode="empty")
        return JsonResponse(
            {"response": "파일 또는 제품 설명을 입력해주세요."},
            status=400,
        )

    return JsonResponse(
        {
            "mode": input_mode,
            "options": options,
            "default_selected_ids": default_selected_ids,
        }
    )


def _parse_selected_summaries(request):
    raw_value = request.POST.get("selectedSummaries", "")
    try:
        values = json.loads(raw_value)
    except (TypeError, json.JSONDecodeError):
        return None
    if not isinstance(values, list):
        return None

    selected = []
    seen = set()
    for value in values:
        summary = preprocess_text(str(value or ""))
        if not summary or summary in seen:
            continue
        if len(summary) > 10000:
            return None
        selected.append(summary)
        seen.add(summary)
    if not 1 <= len(selected) <= 6:
        return None
    return selected


def _parse_search_period(request):
    raw_start = request.POST.get("searchStartDate", "").strip() or "2017-01-01"
    raw_end = request.POST.get("searchEndDate", "").strip()
    try:
        start_date = date.fromisoformat(raw_start)
        end_date = date.fromisoformat(raw_end) if raw_end else None
    except ValueError:
        return None
    if end_date and start_date > end_date:
        return None
    return start_date, end_date


def _search_selected_summaries(request):
    selected_summaries = _parse_selected_summaries(request)
    if not selected_summaries:
        return JsonResponse(
            {"response": "유사도를 판단할 문장을 1개 이상 선택해주세요."},
            status=400,
        )

    search_period = _parse_search_period(request)
    if not search_period:
        return JsonResponse(
            {"response": "인증일자 검색 기간을 올바르게 입력해주세요."},
            status=400,
        )
    cert_date_from, cert_date_to = search_period

    set_request_log_context(
        request,
        feature="similar",
        input_mode=request.POST.get("inputMode", ""),
        search_query=" | ".join(selected_summaries),
    )

    try:
        faiss_result, _ = compare_multiple_from_index(
            selected_summaries,
            k=30,
            cert_date_from=cert_date_from,
            cert_date_to=cert_date_to,
        )
    except SimilarSearchDependencyError as exc:
        return JsonResponse({"response": str(exc)}, status=503)

    rerank_error = ""
    try:
        compare_result = rerank_multiple_similar_candidates(
            selected_summaries,
            faiss_result,
        )
        if not compare_result:
            rerank_error = "LLM 재평가 결과가 비어 있어 FAISS 평균 결과를 표시합니다."
            compare_result = faiss_result
    except (GemmaConfigError, GemmaGenerationError) as exc:
        rerank_error = f"LLM 재평가를 수행하지 못해 FAISS 평균 결과를 표시합니다: {exc}"
        compare_result = faiss_result

    for row in compare_result:
        row.pop("faiss_scores", None)

    similarity_list = [row.get("similarity", 0.0) for row in compare_result]
    set_request_log_context(request, result_count=len(compare_result))
    return JsonResponse(
        {
            "summary": selected_summaries,
            "response": compare_result,
            "similarities": similarity_list,
            "rerank_error": rerank_error,
            "search_period": {
                "start": cert_date_from.isoformat(),
                "end": cert_date_to.isoformat() if cert_date_to else "",
            },
        }
    )


# Django 뷰 함수 (추천 문장 준비 + 선택 문장 검색 API)
@csrf_exempt
def summarize_document(request):
    if request.method != "POST":
        return JsonResponse(
            {"response": "POST 메소드만 지원됩니다."},
            status=405,
        )

    action = request.POST.get("action", "prepare")
    if action == "prepare":
        return _prepare_summary_options(request)
    if action == "search":
        return _search_selected_summaries(request)
    return JsonResponse({"response": "지원하지 않는 요청입니다."}, status=400)
