"""Structured document extraction for the similar-product analysis flow.

Every supported format is converted to small, source-addressable text units.
Legacy Microsoft Office formats are converted with the locally installed
desktop Office applications.  HWP 5 binary files are parsed directly through
their OLE streams; HWPX is parsed as OWPML XML.
"""

from __future__ import annotations

from dataclasses import asdict, dataclass, field
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Iterable
from zipfile import BadZipFile, ZipFile
import contextlib
import os
import re
import shutil
import struct
import subprocess
import sys
import unicodedata
import zlib

import fitz
from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".doc",
    ".docx",
    ".xls",
    ".xlsx",
    ".hwp",
    ".hwpx",
    ".ppt",
    ".pptx",
    ".md",
}

MAX_PDF_PAGES_PER_FILE = 500
MAX_SLIDES_PER_FILE = 500
MAX_EXCEL_CELLS_PER_FILE = 200_000
OCR_TEXT_THRESHOLD = 40
# Tesseract's default page segmentation (3) routes mixed Korean/Latin blocks to the
# wrong script and drops most Hangul.  Mode 4 (single column, variable sizes) reads
# scanned report pages reliably.
OCR_PAGE_SEGMENTATION_MODE = 4
MAX_ARCHIVE_ENTRIES = 20_000
MAX_ARCHIVE_EXPANDED_BYTES = 1024 * 1024 * 1024
MAX_ARCHIVE_COMPRESSION_RATIO = 200


class DocumentParseError(RuntimeError):
    """A safe, user-facing parsing failure."""


@dataclass
class DocumentUnit:
    source_id: str
    filename: str
    kind: str
    locator: str
    text: str
    metadata: dict = field(default_factory=dict)

    def to_dict(self):
        return asdict(self)


@dataclass
class ParsedDocument:
    filename: str
    extension: str
    units: list[DocumentUnit] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    stats: dict = field(default_factory=dict)

    @property
    def text(self):
        return "\n".join(
            f"[{unit.source_id}] {unit.text}" for unit in self.units if unit.text
        )


def normalize_unit_text(value) -> str:
    text = unicodedata.normalize("NFC", str(value or ""))
    text = text.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    text = "".join(char for char in text if char == "\n" or char == "\t" or ord(char) >= 32)
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def _append_unit(
    document: ParsedDocument,
    *,
    kind: str,
    locator: str,
    text,
    metadata: dict | None = None,
):
    normalized = normalize_unit_text(text)
    if not normalized:
        return
    source_id = f"F:{document.filename}|{locator}"
    document.units.append(
        DocumentUnit(
            source_id=source_id,
            filename=document.filename,
            kind=kind,
            locator=locator,
            text=normalized,
            metadata=metadata or {},
        )
    )


def _parse_pdf(path: Path, filename: str) -> ParsedDocument:
    result = ParsedDocument(filename, ".pdf")
    ocr_pages = 0
    with fitz.open(path) as pdf:
        if pdf.needs_pass:
            raise DocumentParseError("암호가 설정된 PDF입니다. 암호를 해제한 파일을 업로드해주세요.")
        page_count = min(len(pdf), MAX_PDF_PAGES_PER_FILE)
        if len(pdf) > page_count:
            result.warnings.append(
                f"PDF {len(pdf)}쪽 중 안전 한도인 {page_count}쪽까지 분석했습니다."
            )
        for page_index in range(page_count):
            page = pdf[page_index]
            blocks = sorted(
                page.get_text("blocks"),
                key=lambda item: (round(item[1] / 8), item[0]),
            )
            page_text = "\n".join(str(block[4]) for block in blocks)
            if len(normalize_unit_text(page_text)) < OCR_TEXT_THRESHOLD:
                ocr_text = _ocr_pdf_page(page)
                if ocr_text is None:
                    warning = "일부 스캔 페이지에 OCR을 적용하지 못했습니다. Tesseract kor/eng 설치를 확인해주세요."
                    if warning not in result.warnings:
                        result.warnings.append(warning)
                elif ocr_text:
                    page_text = ocr_text
                    ocr_pages += 1
            _append_unit(
                result,
                kind="page",
                locator=f"PAGE:{page_index + 1}",
                text=page_text,
                metadata={"page": page_index + 1},
            )
    result.stats.update(page_count=page_count, ocr_pages=ocr_pages)
    return result


def _ocr_pdf_page(page) -> str | None:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return None
    try:
        if not shutil.which("tesseract"):
            installed = Path(os.environ.get("ProgramFiles", r"C:\Program Files")) / "Tesseract-OCR" / "tesseract.exe"
            if installed.exists():
                pytesseract.pytesseract.tesseract_cmd = str(installed)
        local_tessdata = (
            Path(os.environ.get("LOCALAPPDATA", ""))
            / "GSCert"
            / "tessdata"
        )
        if (local_tessdata / "kor.traineddata").exists():
            # pytesseract splits ``config`` with ``shlex.split(posix=False)``, which
            # keeps the quote characters inside the token, so a quoted
            # ``--tessdata-dir`` path reaches Tesseract with the quotes attached and
            # every language fails to load.  Hand the directory over through the
            # environment instead, which also survives paths containing spaces.
            os.environ["TESSDATA_PREFIX"] = str(local_tessdata)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
        return pytesseract.image_to_string(
            image,
            lang="kor+eng",
            config=f"--psm {OCR_PAGE_SEGMENTATION_MODE}",
        )
    except (pytesseract.TesseractNotFoundError, pytesseract.TesseractError):
        return None


def _safe_xml(zip_file: ZipFile, member: str):
    parser = etree.XMLParser(resolve_entities=False, no_network=True, recover=True)
    return etree.fromstring(zip_file.read(member), parser=parser)


def _validate_document_container(path: Path, extension: str):
    with path.open("rb") as source:
        signature = source.read(8)
    if extension == ".pdf" and not signature.startswith(b"%PDF-"):
        raise DocumentParseError("확장자는 PDF이지만 실제 PDF 파일이 아닙니다.")
    if extension in {".doc", ".xls", ".ppt", ".hwp"} and signature != bytes.fromhex(
        "D0CF11E0A1B11AE1"
    ):
        raise DocumentParseError("확장자와 실제 OLE 문서 형식이 일치하지 않습니다.")
    if extension in {".docx", ".xlsx", ".pptx", ".hwpx"}:
        if not signature.startswith(b"PK"):
            raise DocumentParseError("확장자와 실제 OOXML/OWPML 문서 형식이 일치하지 않습니다.")
        try:
            with ZipFile(path) as archive:
                entries = archive.infolist()
                if len(entries) > MAX_ARCHIVE_ENTRIES:
                    raise DocumentParseError("문서 내부 파일 수가 안전 한도를 초과했습니다.")
                expanded = sum(item.file_size for item in entries)
                compressed = sum(max(item.compress_size, 1) for item in entries)
                if expanded > MAX_ARCHIVE_EXPANDED_BYTES:
                    raise DocumentParseError("문서 압축 해제 크기가 안전 한도를 초과했습니다.")
                if expanded / max(compressed, 1) > MAX_ARCHIVE_COMPRESSION_RATIO:
                    raise DocumentParseError("비정상적으로 높은 문서 압축률이 감지되었습니다.")
        except BadZipFile as exc:
            raise DocumentParseError("손상된 압축 문서입니다.") from exc


def _xml_text(root) -> str:
    return normalize_unit_text(" ".join(root.itertext()))


def _parse_docx(path: Path, filename: str) -> ParsedDocument:
    result = ParsedDocument(filename, ".docx")
    try:
        with ZipFile(path) as archive:
            members = set(archive.namelist())
            ordered = ["word/document.xml"]
            ordered.extend(sorted(name for name in members if re.fullmatch(r"word/header\d+\.xml", name)))
            ordered.extend(sorted(name for name in members if re.fullmatch(r"word/footer\d+\.xml", name)))
            ordered.extend(
                name
                for name in ("word/footnotes.xml", "word/endnotes.xml", "word/comments.xml")
                if name in members
            )
            for member in ordered:
                if member not in members:
                    continue
                root = _safe_xml(archive, member)
                blocks = root.xpath(
                    ".//*[local-name()='body']/* | "
                    "./*[local-name()='p' or local-name()='tbl']"
                )
                if not blocks:
                    blocks = [root]
                for index, block in enumerate(blocks, 1):
                    tag = etree.QName(block).localname
                    if tag == "tbl":
                        rows = []
                        for row in block.xpath(".//*[local-name()='tr']"):
                            cells = [
                                _xml_text(cell)
                                for cell in row.xpath("./*[local-name()='tc']")
                            ]
                            if any(cells):
                                rows.append(" | ".join(cell or "-" for cell in cells))
                        text = "\n".join(rows)
                        kind = "table"
                    else:
                        text = _xml_text(block)
                        kind = "paragraph"
                    _append_unit(
                        result,
                        kind=kind,
                        locator=f"{member.upper()}:{index}",
                        text=text,
                        metadata={"part": member},
                    )
    except (BadZipFile, KeyError, etree.XMLSyntaxError) as exc:
        raise DocumentParseError("DOCX 구조를 읽을 수 없습니다.") from exc
    result.stats["unit_count"] = len(result.units)
    return result


def _shape_texts(shape) -> Iterable[tuple[str, str]]:
    if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            yield from _shape_texts(child)
    if getattr(shape, "has_text_frame", False):
        text = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if normalize_unit_text(text):
            yield "text", text
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            rows.append(" | ".join(cell.text.strip() or "-" for cell in row.cells))
        if rows:
            yield "table", "\n".join(rows)
    chart = getattr(shape, "chart", None) if getattr(shape, "has_chart", False) else None
    if chart is not None:
        values = []
        if chart.has_title:
            values.append(chart.chart_title.text_frame.text)
        for series in chart.series:
            values.append(str(getattr(series, "name", "") or ""))
        if any(values):
            yield "chart", "\n".join(values)


def _parse_pptx(path: Path, filename: str) -> ParsedDocument:
    result = ParsedDocument(filename, ".pptx")
    presentation = Presentation(path)
    slide_count = min(len(presentation.slides), MAX_SLIDES_PER_FILE)
    if len(presentation.slides) > slide_count:
        result.warnings.append(
            f"프레젠테이션 {len(presentation.slides)}장 중 {slide_count}장까지 분석했습니다."
        )
    for slide_index, slide in enumerate(presentation.slides, 1):
        if slide_index > slide_count:
            break
        ordered_shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
        for shape_index, shape in enumerate(ordered_shapes, 1):
            for kind, text in _shape_texts(shape):
                _append_unit(
                    result,
                    kind=kind,
                    locator=f"SLIDE:{slide_index}|SHAPE:{shape_index}|{kind.upper()}",
                    text=text,
                    metadata={"slide": slide_index},
                )
        with contextlib.suppress(Exception):
            notes = slide.notes_slide.notes_text_frame.text
            _append_unit(
                result,
                kind="notes",
                locator=f"SLIDE:{slide_index}|NOTES",
                text=notes,
                metadata={"slide": slide_index},
            )
    result.stats["slide_count"] = slide_count
    return result


def _parse_xlsx(path: Path, filename: str) -> ParsedDocument:
    result = ParsedDocument(filename, ".xlsx")
    workbook_values = load_workbook(path, read_only=True, data_only=True)
    workbook_formulas = load_workbook(path, read_only=True, data_only=False)
    cell_count = 0
    try:
        for value_sheet, formula_sheet in zip(
            workbook_values.worksheets, workbook_formulas.worksheets
        ):
            for row_index, (value_row, formula_row) in enumerate(
                zip(
                    value_sheet.iter_rows(),
                    formula_sheet.iter_rows(),
                ),
                1,
            ):
                cells = []
                first_col = None
                last_col = None
                for value_cell, formula_cell in zip(value_row, formula_row):
                    value = value_cell.value
                    formula = formula_cell.value
                    if value is None and formula is None:
                        continue
                    cell_count += 1
                    if cell_count > MAX_EXCEL_CELLS_PER_FILE:
                        result.warnings.append(
                            f"비어 있지 않은 셀 {MAX_EXCEL_CELLS_PER_FILE:,}개까지 분석했습니다."
                        )
                        break
                    first_col = first_col or value_cell.column_letter
                    last_col = value_cell.column_letter
                    display = normalize_unit_text(value if value is not None else formula)
                    if isinstance(formula, str) and formula.startswith("="):
                        display = f"{display} (수식: {formula})" if value is not None else formula
                    cells.append(f"{value_cell.coordinate}={display}")
                if cells:
                    _append_unit(
                        result,
                        kind="table_row",
                        locator=f"SHEET:{value_sheet.title}|ROW:{row_index}",
                        text=" | ".join(cells),
                        metadata={
                            "sheet": value_sheet.title,
                            "row": row_index,
                            "range": f"{first_col}{row_index}:{last_col}{row_index}",
                        },
                    )
                if cell_count > MAX_EXCEL_CELLS_PER_FILE:
                    break
            if cell_count > MAX_EXCEL_CELLS_PER_FILE:
                break
    finally:
        workbook_values.close()
        workbook_formulas.close()
    result.stats["non_empty_cells"] = min(cell_count, MAX_EXCEL_CELLS_PER_FILE)
    return result


def _parse_md(path: Path, filename: str) -> ParsedDocument:
    result = ParsedDocument(filename, ".md")
    raw = None
    for encoding in ("utf-8-sig", "utf-8", "cp949"):
        try:
            raw = path.read_text(encoding=encoding)
            break
        except UnicodeDecodeError:
            continue
    if raw is None:
        raise DocumentParseError("Markdown 파일의 문자 인코딩을 확인해주세요.")

    current_heading = "문서 시작"
    block = []
    in_code = False

    def flush():
        nonlocal block
        text = "\n".join(block)
        if text.strip():
            _append_unit(
                result,
                kind="section",
                locator=f"SECTION:{len(result.units) + 1}|{current_heading[:80]}",
                text=text,
                metadata={"heading": current_heading},
            )
        block = []

    for line in raw.splitlines():
        if line.lstrip().startswith("```"):
            in_code = not in_code
            block.append(line)
            continue
        heading = re.match(r"^\s{0,3}#{1,6}\s+(.+?)\s*$", line)
        if heading and not in_code:
            flush()
            current_heading = heading.group(1)
            block.append(line)
        else:
            block.append(line)
    flush()
    return result


def _parse_hwpx(path: Path, filename: str) -> ParsedDocument:
    result = ParsedDocument(filename, ".hwpx")
    try:
        with ZipFile(path) as archive:
            sections = sorted(
                name
                for name in archive.namelist()
                if re.fullmatch(r"Contents/section\d+\.xml", name, re.IGNORECASE)
            )
            if not sections:
                raise DocumentParseError("HWPX 본문 섹션을 찾을 수 없습니다.")
            for section_index, member in enumerate(sections, 1):
                root = _safe_xml(archive, member)
                paragraphs = root.xpath(".//*[local-name()='p']")
                for paragraph_index, paragraph in enumerate(paragraphs, 1):
                    _append_unit(
                        result,
                        kind="paragraph",
                        locator=f"SECTION:{section_index}|PARAGRAPH:{paragraph_index}",
                        text=_xml_text(paragraph),
                        metadata={"section": section_index},
                    )
    except (BadZipFile, etree.XMLSyntaxError) as exc:
        raise DocumentParseError("HWPX 구조를 읽을 수 없습니다.") from exc
    return result


def _decode_hwp_para_text(payload: bytes) -> str:
    text = payload.decode("utf-16le", errors="ignore")
    # HWP inline controls occupy one or several UTF-16 code units.  Keeping
    # tabs/newlines and replacing the rest with a space preserves word breaks.
    cleaned = []
    index = 0
    while index < len(text):
        code = ord(text[index])
        if code in (9, 10, 13):
            cleaned.append("\n" if code in (10, 13) else "\t")
        elif code < 32:
            cleaned.append(" ")
        else:
            cleaned.append(text[index])
        index += 1
    return "".join(cleaned)


def _parse_hwp(path: Path, filename: str) -> ParsedDocument:
    try:
        import olefile
    except ImportError as exc:
        raise DocumentParseError(
            "HWP 파서(olefile)가 설치되지 않았습니다. requirements.txt를 설치해주세요."
        ) from exc

    result = ParsedDocument(filename, ".hwp")
    if not olefile.isOleFile(str(path)):
        raise DocumentParseError("유효한 HWP 5.x OLE 문서가 아닙니다.")
    with olefile.OleFileIO(str(path)) as ole:
        if not ole.exists("FileHeader"):
            raise DocumentParseError("HWP FileHeader를 찾을 수 없습니다.")
        header = ole.openstream("FileHeader").read()
        compressed = bool(header[36] & 0x01) if len(header) > 36 else False
        sections = sorted(
            entry
            for entry in ole.listdir()
            if len(entry) == 2 and entry[0] == "BodyText" and entry[1].startswith("Section")
        )
        for section_index, stream_path in enumerate(sections, 1):
            stream = ole.openstream(stream_path).read()
            if compressed:
                try:
                    stream = zlib.decompress(stream, -15)
                except zlib.error as exc:
                    raise DocumentParseError("압축된 HWP 본문을 해제하지 못했습니다.") from exc
            offset = 0
            paragraph_index = 0
            while offset + 4 <= len(stream):
                header_value = struct.unpack_from("<I", stream, offset)[0]
                offset += 4
                tag_id = header_value & 0x3FF
                size = (header_value >> 20) & 0xFFF
                if size == 0xFFF:
                    if offset + 4 > len(stream):
                        break
                    size = struct.unpack_from("<I", stream, offset)[0]
                    offset += 4
                payload = stream[offset : offset + size]
                offset += size
                if tag_id == 67:  # HWPTAG_PARA_TEXT
                    paragraph_index += 1
                    _append_unit(
                        result,
                        kind="paragraph",
                        locator=f"SECTION:{section_index}|PARAGRAPH:{paragraph_index}",
                        text=_decode_hwp_para_text(payload),
                        metadata={"section": section_index},
                    )
    if not result.units:
        raise DocumentParseError("HWP 본문 텍스트를 추출하지 못했습니다.")
    return result


def _convert_legacy_office(path: Path, extension: str, output_dir: Path) -> Path:
    destination_extensions = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}
    converted_extension = destination_extensions.get(extension)
    if not converted_extension:
        raise DocumentParseError("지원하지 않는 Office 변환 형식입니다.")
    destination = output_dir / f"{path.stem}{converted_extension}"
    creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
    try:
        completed = subprocess.run(
            [
                sys.executable,
                "-m",
                "main.utils.office_convert",
                str(path.resolve()),
                str(destination.resolve()),
            ],
            capture_output=True,
            text=True,
            timeout=180,
            creationflags=creation_flags,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise DocumentParseError(
            f"MS Office {extension} 변환이 3분 제한 시간을 초과했습니다."
        ) from exc
    if completed.returncode != 0 or not destination.exists():
        raise DocumentParseError(
            f"MS Office에서 {extension} 파일을 변환하지 못했습니다."
        )
    return destination


def parse_document(path, original_name: str | None = None) -> ParsedDocument:
    source_path = Path(path)
    filename = Path(original_name or source_path.name).name
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise DocumentParseError(
            "지원 형식은 pdf, doc(x), xls(x), hwp(x), ppt(x), md입니다."
        )
    _validate_document_container(source_path, extension)
    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".hwp": _parse_hwp,
        ".hwpx": _parse_hwpx,
        ".pptx": _parse_pptx,
        ".md": _parse_md,
    }
    if extension in (".doc", ".xls", ".ppt"):
        with TemporaryDirectory(prefix="similar-office-") as temp_dir:
            converted = _convert_legacy_office(source_path, extension, Path(temp_dir))
            converted_extension = converted.suffix.lower()
            parsed = parsers[converted_extension](converted, filename)
            parsed.extension = extension
            parsed.warnings.insert(
                0, f"{extension} 파일을 MS Office로 {converted_extension} 형식으로 변환해 분석했습니다."
            )
            return parsed
    return parsers[extension](source_path, filename)


def deduplicate_units(documents: Iterable[ParsedDocument]):
    unique = []
    duplicate_count = 0
    seen = set()
    for document in documents:
        for unit in document.units:
            key = re.sub(r"\W+", "", unit.text).casefold()
            if not key:
                key = unit.text.casefold().strip()
            if not key:
                continue
            if key in seen:
                duplicate_count += 1
                continue
            seen.add(key)
            unique.append(unit)
    return unique, duplicate_count


def save_uploaded_file(uploaded_file, destination: Path):
    destination.parent.mkdir(parents=True, exist_ok=True)
    with destination.open("wb") as output:
        for chunk in uploaded_file.chunks():
            output.write(chunk)
